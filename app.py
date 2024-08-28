from flask import Flask, request, jsonify
from pptx import Presentation
from io import BytesIO
import os
import requests
import sys
import shutil
from pdf2image import convert_from_path

app = Flask(__name__)

# 標準出力と標準エラー出力のバッファリングを無効にする
sys.stdout = open(sys.stdout.fileno(), mode='w', buffering=1)
sys.stderr = open(sys.stderr.fileno(), mode='w', buffering=1)

def download_pptx(url):
    response = requests.get(url)
    response.raise_for_status()
    return BytesIO(response.content)

def save_presentation_as_pdf(pptx_bytes, pdf_path):
    # Save the PowerPoint presentation as a temporary PPTX file
    temp_pptx_path = "temp.pptx"
    with open(temp_pptx_path, "wb") as f:
        f.write(pptx_bytes.getbuffer())
    
    # Convert the PPTX file to PDF using unoconv
    os.system(f"unoconv -f pdf {temp_pptx_path}")
    
    # Move the generated PDF to the desired location
    temp_pdf_path = temp_pptx_path.replace(".pptx", ".pdf")
    shutil.move(temp_pdf_path, pdf_path)
    
    # Remove the temporary PPTX file
    os.remove(temp_pptx_path)

def extract_slides(pptx_bytes, output_dir):
    presentation = Presentation(pptx_bytes)
    slides = []
    text_output = []

    # Save presentation as PDF
    pdf_path = os.path.join(output_dir, "presentation.pdf")
    save_presentation_as_pdf(pptx_bytes, pdf_path)

    # Convert PDF to images
    images = convert_from_path(pdf_path, dpi=300)

    for i, slide in enumerate(presentation.slides):
        slide_text = []
        
        # Extract text
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        # Combine text and add to output
        slide_text_combined = "\n".join(slide_text)
        text_output.append(f"Slide {i+1}:\n{slide_text_combined}\n")
        
        # Save slide as PNG
        slide_filename = os.path.join(output_dir, f"slide_{i+1}.png")
        images[i].save(slide_filename, 'PNG')
    
    # Save all text to a single file
    text_filename = os.path.join(output_dir, "slides_text.txt")
    with open(text_filename, "w") as text_file:
        text_file.write("\n".join(text_output))

    return text_filename

@app.route('/analyze', methods=['POST'])
def analyze():
    url = request.json.get("url")
    output_dir = "/output"
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    pptx_bytes = download_pptx(url)
    text_filename = extract_slides(pptx_bytes, output_dir)
    
    # Print extracted text to standard output
    with open(text_filename, "r") as text_file:
        print(text_file.read())
    
    return jsonify({"message": "Analysis complete", "text_file": text_filename})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000)